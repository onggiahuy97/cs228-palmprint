"""Generate CS 228 project proposal presentation as .pptx using python-pptx."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent
FIGURES = ROOT / "figures"
CHECKPOINTS = ROOT / "checkpoints" / "baseline_v2"
OUTPUT = ROOT / "docs" / "proposal_presentation.pptx"

# Colors
BG_COLOR = RGBColor(0x1A, 0x1A, 0x2E)
TITLE_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
BODY_COLOR = RGBColor(0xDD, 0xDD, 0xDD)
ACCENT_COLOR = RGBColor(0x4E, 0xC9, 0xB0)
SUBTLE_COLOR = RGBColor(0x99, 0x99, 0x99)
RED_ACCENT = RGBColor(0xFF, 0x6B, 0x6B)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color=BG_COLOR):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=BODY_COLOR, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=18,
                    color=BODY_COLOR, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing

        # Support rich text: list of (text, kwargs) tuples or plain string
        if isinstance(item, list):
            for segment in item:
                if isinstance(segment, tuple):
                    txt, kwargs = segment
                    run = p.add_run()
                    run.text = txt
                    run.font.size = Pt(kwargs.get("size", font_size))
                    run.font.color.rgb = kwargs.get("color", color)
                    run.font.bold = kwargs.get("bold", False)
                    run.font.name = kwargs.get("font", "Calibri")
                else:
                    run = p.add_run()
                    run.text = segment
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = color
                    run.font.name = "Calibri"
        else:
            p.text = f"  \u2022  {item}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.name = "Calibri"
    return tf


def add_slide_title(slide, title, subtitle=None):
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8),
                 title, font_size=32, color=TITLE_COLOR, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11.5), Inches(0.5),
                     subtitle, font_size=16, color=SUBTLE_COLOR)


def add_image_centered(slide, img_path, top=Inches(2.0), max_w=Inches(10), max_h=Inches(4.8)):
    from PIL import Image
    img = Image.open(img_path)
    w, h = img.size
    aspect = w / h
    # Fit within max bounds
    if aspect > (max_w / max_h):
        width = max_w
        height = int(width / aspect)
    else:
        height = max_h
        width = int(height * aspect)
    left = (SLIDE_WIDTH - width) // 2
    slide.shapes.add_picture(str(img_path), left, top, width, height)


def add_table(slide, left, top, width, height, rows, col_widths=None,
              header=True, font_size=14):
    n_rows = len(rows)
    n_cols = len(rows[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r, row in enumerate(rows):
        for c, cell_text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "Calibri"
                if r == 0 and header:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                else:
                    paragraph.font.color.rgb = BODY_COLOR
            # Cell background
            cell_fill = cell.fill
            cell_fill.solid()
            if r == 0 and header:
                cell_fill.fore_color.rgb = RGBColor(0x2D, 0x2D, 0x44)
            elif r % 2 == 0:
                cell_fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3A)
            else:
                cell_fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    return table


def add_speaker_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


# ──────────────────────────────────────────────
# SLIDES
# ──────────────────────────────────────────────

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)
    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
                 "Robust Contactless Palmprint Verification",
                 font_size=40, color=TITLE_COLOR, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
                 "with Deep Metric Learning & Robustness Benchmarking",
                 font_size=24, color=ACCENT_COLOR, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
                 "Huy Ong  \u2022  CS 228: Biometric Security with AI",
                 font_size=18, color=SUBTLE_COLOR, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.5),
                 "Project Proposal  \u2022  Spring 2026",
                 font_size=16, color=SUBTLE_COLOR, alignment=PP_ALIGN.CENTER)
    add_speaker_notes(slide, "Introduce yourself. Solo project. ~15 min presentation.")


def slide_motivation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "Motivation")
    items = [
        "Contactless palmprint: hygienic, user-friendly, rich features (lines, wrinkles, texture)",
        "No physical contact \u2192 growing demand post-COVID",
        "But contactless capture = noisy: pose, lighting, blur, distance all vary",
        "Key question: How badly do these real-world factors break verification?",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(4.5),
                    items, font_size=22)
    add_speaker_notes(slide,
        "Hook: 'What happens when capture conditions aren't perfect?' "
        "Palmprints are attractive because they're contactless and distinctive, "
        "but that same contactless capture introduces noise. This project asks: how bad is it?")


def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "Problem Statement")
    items = [
        "1.  Build a deep-learning palmprint verifier (ResNet18 + ArcFace)",
        "2.  Stress-test it under 9 corruption types \u00d7 5 severity levels",
        "3.  Identify the most damaging real-world factors",
        "4.  Propose and implement mitigations to improve robustness",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(4.5),
                    items, font_size=22)
    add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.6),
                 "Security angle: FAR matters \u2014 a degraded image shouldn't let impostors through.",
                 font_size=16, color=RED_ACCENT)
    add_speaker_notes(slide,
        "Emphasize the security angle. In biometrics, it's not just about convenience \u2014 "
        "False Accept Rate is a security metric. We need to know how corruptions affect FAR.")


def slide_lit_survey(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "Literature \u2014 Contactless Palmprint Survey",
                    "Fei et al., IEEE T-BIOM 2021")
    items = [
        "Comprehensive survey of contactless palmprint recognition pipeline",
        "Covers: ROI extraction \u2192 feature extraction \u2192 matching",
        "Catalogs open problems: lighting, pose, partial occlusion",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.8), Inches(7), Inches(2.5),
                    items, font_size=20)
    add_text_box(slide, Inches(0.8), Inches(4.5), Inches(5), Inches(0.5),
                 "Strengths:", font_size=18, color=ACCENT_COLOR, bold=True)
    add_bullet_list(slide, Inches(0.8), Inches(5.0), Inches(5.5), Inches(1.2),
                    ["Broad coverage of the full pipeline",
                     "Identifies key open challenges"],
                    font_size=16, color=BODY_COLOR)
    add_text_box(slide, Inches(6.5), Inches(4.5), Inches(5), Inches(0.5),
                 "Gaps:", font_size=18, color=RED_ACCENT, bold=True)
    add_bullet_list(slide, Inches(6.5), Inches(5.0), Inches(5.5), Inches(1.2),
                    ["Qualitative \u2014 no quantitative robustness eval",
                     "No standardized corruption benchmark"],
                    font_size=16, color=BODY_COLOR)
    add_speaker_notes(slide,
        "This survey paper gives a great overview of the field but stays qualitative. "
        "It identifies problems like lighting and pose but doesn't measure how bad they are.")


def slide_lit_deep(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "Literature \u2014 Deep Features for Palmprint",
                    "Fei et al., IEEE TIFS 2020")
    items = [
        "ResNet backbone + ArcFace angular margin loss",
        "Achieves <1.5% EER on contactless palmprint datasets",
        "Demonstrates that deep metric learning outperforms handcrafted features",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.8), Inches(7), Inches(2.5),
                    items, font_size=20)
    add_text_box(slide, Inches(0.8), Inches(4.5), Inches(5), Inches(0.5),
                 "Strengths:", font_size=18, color=ACCENT_COLOR, bold=True)
    add_bullet_list(slide, Inches(0.8), Inches(5.0), Inches(5.5), Inches(1.2),
                    ["Strong empirical results",
                     "Clean, reproducible methodology"],
                    font_size=16, color=BODY_COLOR)
    add_text_box(slide, Inches(6.5), Inches(4.5), Inches(5), Inches(0.5),
                 "Gaps:", font_size=18, color=RED_ACCENT, bold=True)
    add_bullet_list(slide, Inches(6.5), Inches(5.0), Inches(5.5), Inches(1.2),
                    ["No robustness evaluation under corruptions",
                     "Only tested on clean capture conditions"],
                    font_size=16, color=BODY_COLOR)
    add_speaker_notes(slide,
        "This paper is the closest to our approach \u2014 same backbone family and loss. "
        "Key gap: they only evaluate on clean data. We extend this with robustness testing.")


def slide_lit_corruption(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "Literature \u2014 Corruption Benchmark",
                    "Hendrycks & Dietterich, ICLR 2019")
    items = [
        "Introduced ImageNet-C: systematic corruption benchmark for classifiers",
        "15 corruption types at 5 severity levels",
        "Showed that even top classifiers degrade significantly under common corruptions",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.8), Inches(7), Inches(2.5),
                    items, font_size=20)
    add_text_box(slide, Inches(0.8), Inches(4.5), Inches(5), Inches(0.5),
                 "Strengths:", font_size=18, color=ACCENT_COLOR, bold=True)
    add_bullet_list(slide, Inches(0.8), Inches(5.0), Inches(5.5), Inches(1.2),
                    ["Standardized benchmark methodology",
                     "Widely adopted in robustness research"],
                    font_size=16, color=BODY_COLOR)
    add_text_box(slide, Inches(6.5), Inches(4.5), Inches(5), Inches(0.5),
                 "Gaps:", font_size=18, color=RED_ACCENT, bold=True)
    add_bullet_list(slide, Inches(6.5), Inches(5.0), Inches(5.5), Inches(1.2),
                    ["Classification-only \u2014 no biometric metrics (EER, FAR)",
                     "Corruptions not tailored to palmprint capture"],
                    font_size=16, color=BODY_COLOR)
    add_speaker_notes(slide,
        "We adapt Hendrycks' paradigm from classification to biometric verification. "
        "Instead of top-1 accuracy, we measure EER. Instead of generic corruptions, "
        "we use palmprint-specific ones like hand rotation and occlusion.")


def slide_approach_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "Our Approach \u2014 Overview")
    # Pipeline boxes
    modules = [
        ("Dataset &\nSplits", ACCENT_COLOR),
        ("Baseline\nVerifier", ACCENT_COLOR),
        ("Corruption\nBenchmark", ACCENT_COLOR),
        ("Mitigations", RGBColor(0xFF, 0xA5, 0x00)),
        ("Analytics &\nReporting", ACCENT_COLOR),
    ]
    box_w = Inches(2.0)
    box_h = Inches(1.2)
    gap = Inches(0.3)
    total_w = len(modules) * box_w + (len(modules) - 1) * gap
    start_x = (SLIDE_WIDTH - total_w) // 2
    top_y = Inches(2.5)

    for i, (label, color) in enumerate(modules):
        left = start_x + i * (box_w + gap)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       left, top_y, box_w, box_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x2D, 0x2D, 0x44)
        shape.line.color.rgb = color
        shape.line.width = Pt(2)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(16)
        p.font.color.rgb = TITLE_COLOR
        p.font.bold = True
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Arrow between boxes
        if i < len(modules) - 1:
            arrow_left = left + box_w
            arrow_top = top_y + box_h // 2 - Inches(0.1)
            add_text_box(slide, arrow_left, arrow_top, gap, Inches(0.3),
                         "\u2192", font_size=24, color=SUBTLE_COLOR,
                         alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.8), Inches(4.5), Inches(11), Inches(2),
                 "5-module architecture: from raw data to deployment recommendations.\n"
                 "Orange = proposed future work (mitigations).",
                 font_size=16, color=SUBTLE_COLOR)
    add_speaker_notes(slide,
        "High-level architecture. The first 3 modules are already implemented. "
        "Module 4 (mitigations) is the proposed remaining work.")


def slide_dataset(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "Module 1: Dataset & Splits",
                    "PolyU\u2013IITD Contactless Palmprint v3.0")
    items = [
        "611 subjects, 20 images each (10 left + 10 right hand)",
        "12,220 total images \u2192 128\u00d7128 grayscale ROI crops",
        "Each (subject, hand) = unique identity \u2192 1,222 classes",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.8), Inches(5), Inches(2),
                    items, font_size=20)

    rows = [
        ["Split", "Subjects", "Images", "Classes"],
        ["Train (70%)", "427", "8,540", "854"],
        ["Val (10%)", "61", "1,220", "122"],
        ["Test (20%)", "123", "2,460", "246"],
    ]
    add_table(slide, Inches(6.5), Inches(1.8), Inches(5.5), Inches(2.5),
              rows, font_size=16)

    add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.6),
                 "Subject-disjoint: no identity appears in more than one split \u2192 no leakage.",
                 font_size=16, color=ACCENT_COLOR, bold=True)
    add_speaker_notes(slide,
        "Stress 'no identity leakage' \u2014 this is critical for honest evaluation. "
        "Many papers accidentally leak identities across train/test.")


def slide_baseline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "Module 2: Baseline Verifier",
                    "ResNet18 + ArcFace Angular Margin Loss")
    items = [
        "Backbone: ResNet18 (ImageNet pretrained) \u2192 512-D features",
        "Projection: MLP (512 \u2192 512 \u2192 256) with BN + ReLU",
        "Output: 256-D L2-normalized embeddings on unit hypersphere",
        "Training loss: ArcFace (margin=0.3, scale=30)",
        "Optimizer: AdamW (lr=3e-4, weight_decay=1e-3)",
        "Verification: cosine similarity between embedding pairs",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(4),
                    items, font_size=20)
    add_speaker_notes(slide,
        "Standard deep metric learning setup. ArcFace enforces angular margin "
        "on the hypersphere, producing discriminative embeddings. "
        "The ArcFace head is discarded at inference \u2014 only the embedder is used.")


def slide_training_curves(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "Preliminary Result \u2014 Training Convergence")
    img_path = FIGURES / "training_curves.png"
    add_image_centered(slide, img_path, top=Inches(1.6), max_h=Inches(4.5))
    add_text_box(slide, Inches(0.8), Inches(6.3), Inches(11), Inches(0.5),
                 "Best validation EER: 1.13% at epoch 32. Early stopping applied.",
                 font_size=16, color=ACCENT_COLOR)
    add_speaker_notes(slide,
        "Walk through the 3 panels: loss decreasing, accuracy rising, validation EER "
        "reaching minimum at epoch 32 then rising slightly = mild overfitting. "
        "We use the epoch-32 checkpoint for all evaluation.")


def slide_clean_eer(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "Preliminary Result \u2014 Clean Test Performance")

    # Big EER number
    add_text_box(slide, Inches(0.8), Inches(1.8), Inches(4), Inches(1.2),
                 "1.67%", font_size=72, color=ACCENT_COLOR, bold=True)
    add_text_box(slide, Inches(0.8), Inches(3.0), Inches(4), Inches(0.5),
                 "Equal Error Rate (test set)", font_size=20, color=SUBTLE_COLOR)

    # Metrics table
    rows = [
        ["Metric", "Value"],
        ["EER", "1.67%"],
        ["FAR @ 1% FRR", "4.12%"],
        ["Genuine mean sim.", "0.759 \u00b1 0.160"],
        ["Impostor mean sim.", "0.033 \u00b1 0.105"],
    ]
    add_table(slide, Inches(5.5), Inches(1.8), Inches(6.5), Inches(3),
              rows, font_size=16)

    add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.8),
                 "Our baseline already works \u2014 now let's break it.",
                 font_size=22, color=RED_ACCENT, bold=True)
    add_speaker_notes(slide,
        "1.67% EER is competitive with published results. "
        "Clear separation between genuine (0.76) and impostor (0.03) scores. "
        "Transition: 'The baseline works well on clean data. But what about real-world conditions?'")


def slide_score_dist_roc(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "Score Distribution & ROC Curve")

    # Two images side by side
    from PIL import Image
    left_img = CHECKPOINTS / "score_dist_test.png"
    right_img = FIGURES / "roc_curve.png"
    max_h = Inches(4.2)
    max_w = Inches(5.5)
    top = Inches(2.0)

    for img_path, x_offset in [(left_img, Inches(0.5)), (right_img, Inches(6.8))]:
        img = Image.open(img_path)
        w, h = img.size
        aspect = w / h
        if aspect > (max_w / max_h):
            width = max_w
            height = int(width / aspect)
        else:
            height = max_h
            width = int(height * aspect)
        slide.shapes.add_picture(str(img_path), x_offset, top, width, height)

    add_speaker_notes(slide,
        "Left: score distributions \u2014 genuine pairs (green) cluster around 0.76, "
        "impostors near 0.03. Clear separation. "
        "Right: ROC curve hugs the upper-left corner, AUC > 0.99.")


def slide_embeddings(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "Preliminary Result \u2014 Embedding Space")
    img_path = FIGURES / "tsne_embeddings.png"
    add_image_centered(slide, img_path, top=Inches(1.6), max_h=Inches(4.8))
    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
                 "t-SNE of 20 test identities. Compact, well-separated clusters = good embeddings.",
                 font_size=16, color=ACCENT_COLOR)
    add_speaker_notes(slide,
        "Each color = one identity. Compact clusters with clear separation. "
        "Qualitative confirmation that ArcFace structures the embedding space well.")


def slide_corruption_suite(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "Module 3: Corruption Suite",
                    "9 corruption types \u00d7 5 severity levels = 45 conditions")
    img_path = FIGURES / "corruption_samples.png"
    add_image_centered(slide, img_path, top=Inches(1.6), max_h=Inches(4.5))
    add_text_box(slide, Inches(0.8), Inches(6.3), Inches(11), Inches(0.5),
                 "Each corruption simulates a real-world capture condition (shown at severity 3).",
                 font_size=16, color=SUBTLE_COLOR)
    add_speaker_notes(slide,
        "Walk through each corruption and its real-world analog: "
        "rotation = hand pose, scale = distance, brightness = lighting, "
        "motion blur = hand tremor, noise = sensor/low-light, "
        "occlusion = fingers covering palm, JPEG = compression.")


def slide_heatmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "Robustness Results \u2014 EER Heatmap")
    img_path = CHECKPOINTS / "robustness" / "eer_heatmap.png"
    add_image_centered(slide, img_path, top=Inches(1.6), max_h=Inches(5.0))
    add_text_box(slide, Inches(0.8), Inches(6.7), Inches(11), Inches(0.5),
                 "Darker = higher EER. Clear red zones: Noise, Motion Blur, Brightness.",
                 font_size=16, color=RED_ACCENT)
    add_speaker_notes(slide,
        "Highlight the red zones. Noise at severity 5 = 21.9% EER (13x worse than clean). "
        "Motion blur and brightness also severe. Contrast and scale are relatively safe.")


def slide_eer_curves(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "Robustness Results \u2014 EER vs. Severity")

    img_path = CHECKPOINTS / "robustness" / "eer_curves.png"
    add_image_centered(slide, img_path, top=Inches(1.5), max_h=Inches(3.5))

    rows = [
        ["Corruption", "Clean EER", "Worst EER", "Degradation"],
        ["Noise", "1.73%", "21.90%", "+1,170%"],
        ["Motion Blur", "1.73%", "19.67%", "+1,041%"],
        ["Brightness", "1.73%", "15.45%", "+796%"],
    ]
    add_table(slide, Inches(2.5), Inches(5.3), Inches(8), Inches(1.8),
              rows, font_size=15)
    add_speaker_notes(slide,
        "Top 3 worst corruptions. Noise destroys fine texture details. "
        "Motion blur smears principal lines. Brightness collapses feature visibility. "
        "These are the corruptions we need to mitigate.")


def slide_mitigations(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "Module 4: Proposed Mitigations")

    # Strategy 1
    add_text_box(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.5),
                 "Strategy 1: Robust Augmentation Training",
                 font_size=22, color=ACCENT_COLOR, bold=True)
    items1 = [
        "Add corruption-specific augmentations during training",
        "Focus on top-3 damaging corruptions: noise, motion blur, brightness",
        "Hypothesis: model learns corruption-invariant features",
        "Risk: may slightly hurt clean performance (tradeoff)",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(2.5),
                    items1, font_size=16)

    # Strategy 2
    add_text_box(slide, Inches(6.8), Inches(1.8), Inches(5.5), Inches(0.5),
                 "Strategy 2: Quality Gating",
                 font_size=22, color=ACCENT_COLOR, bold=True)
    items2 = [
        "Estimate image quality before matching",
        "Reject/flag samples below quality threshold",
        "Metrics: blur score, brightness, noise level",
        "Tradeoff: higher security but may increase FRR",
    ]
    add_bullet_list(slide, Inches(6.8), Inches(2.5), Inches(5.5), Inches(2.5),
                    items2, font_size=16)

    add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(0.8),
                 "Will implement at least one strategy and compare baseline vs. improved "
                 "on the full 45-condition benchmark.",
                 font_size=16, color=SUBTLE_COLOR)
    add_speaker_notes(slide,
        "Two complementary strategies. Augmentation training is proactive (makes the model better). "
        "Quality gating is reactive (rejects bad inputs). Will implement at least one and evaluate.")


def slide_timeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "Timeline")

    rows = [
        ["Week", "Milestone", "Status"],
        ["1\u20132", "Dataset splits + training pipeline", "\u2705 Done"],
        ["3\u20134", "Clean verification evaluation + baseline metrics", "\u2705 Done"],
        ["5\u20136", "Corruption suite + robustness benchmark", "\u2705 Done"],
        ["7\u20139", "Implement mitigation (augmentation or quality gate)", "Planned"],
        ["10\u201311", "Re-run full benchmark + finalize results", "Planned"],
        ["12\u201315", "Final report + demo + presentation", "Planned"],
    ]
    add_table(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(4.5),
              rows, font_size=17)

    add_text_box(slide, Inches(0.8), Inches(6.3), Inches(11), Inches(0.5),
                 "Baseline + benchmark already complete. Remaining: mitigations + final evaluation.",
                 font_size=16, color=ACCENT_COLOR, bold=True)
    add_speaker_notes(slide,
        "Highlight that the first 3 milestones are already done. "
        "We're ahead of schedule. Remaining work: mitigations and final write-up.")


def slide_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "Summary & Questions")
    items = [
        "Baseline verifier achieves 1.67% EER on clean data",
        "Systematic robustness benchmark reveals critical vulnerabilities:",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(1.5),
                    items, font_size=22)

    sub_items = [
        "Noise (+1,170%), Motion Blur (+1,041%), Brightness (+796%)",
    ]
    add_bullet_list(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(0.8),
                    sub_items, font_size=20, color=RED_ACCENT)

    items2 = [
        "Two mitigation strategies proposed: augmentation training + quality gating",
        "Next: implement mitigations and demonstrate measurable improvement",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(4.5), Inches(11), Inches(1.5),
                    items2, font_size=22)

    add_text_box(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.8),
                 "Questions?",
                 font_size=36, color=ACCENT_COLOR, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_speaker_notes(slide,
        "Recap the three key points: baseline works, vulnerabilities identified, "
        "mitigations planned. Open for Q&A.")


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Generate all slides
    slide_title(prs)
    slide_motivation(prs)
    slide_problem(prs)
    slide_lit_survey(prs)
    slide_lit_deep(prs)
    slide_lit_corruption(prs)
    slide_approach_overview(prs)
    slide_dataset(prs)
    slide_baseline(prs)
    slide_training_curves(prs)
    slide_clean_eer(prs)
    slide_score_dist_roc(prs)
    slide_embeddings(prs)
    slide_corruption_suite(prs)
    slide_heatmap(prs)
    slide_eer_curves(prs)
    slide_mitigations(prs)
    slide_timeline(prs)
    slide_summary(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Presentation saved to {OUTPUT}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
